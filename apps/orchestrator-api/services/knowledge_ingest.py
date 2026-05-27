"""
knowledge_ingest.py
===================

Assistant Knowledge Base — ingestion + retrieval.

Pipeline:
  1) Parse uploaded file by MIME → list of (location, title, text) chunks
       - xlsx/xls   → openpyxl, one chunk per N-row block, headers prepended
       - csv        → split lines, treat header row + N-row blocks like xlsx
       - pdf        → pypdf, one chunk per page
       - docx       → python-docx, paragraph-grouped chunks
       - pptx       → python-pptx, one chunk per slide
       - txt/md     → plain-text token-based chunks
  2) Embed each chunk via OpenAI text-embedding-3-small (1536-dim).
  3) Insert into Supabase pgvector table (assistant_knowledge_chunks).

Retrieval entry point:
  - rag_retrieve(agent_id, query, top_k=8) → list of dicts with content+location+similarity.
    Used by both:
      - the assistant's tool `search_knowledge_base` (LLM may call it explicitly)
      - the RAG-first shim in assistant_agent.py (auto-prepends top hits to context)
"""

from __future__ import annotations

import io
import os
import re
import json
import uuid
import logging
from typing import Iterable, Optional

import httpx
from sqlalchemy import text as sa_text
from sqlalchemy.orm import Session

log = logging.getLogger("knowledge_ingest")

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------
#
# Embedding provider — set EMBED_PROVIDER=openai or =gemini in env.
# Defaults to gemini (free tier, no quota hassle, 768-dim).
# Both produce dense vectors compatible with pgvector; the DB schema's
# vector dimension MUST match the active provider — see assistant_knowledge
# migrations.

# Default to "none" — store chunks without embeddings, retrieve via pgroonga
# full-text search. This is the path that works on Render's free tier (no
# fastembed install, no model download, ~0 RAM overhead). Set EMBED_PROVIDER
# to 'openai' / 'gemini' / 'local' to additionally generate vectors.
EMBED_PROVIDER = (os.getenv("EMBED_PROVIDER") or "none").lower()

if EMBED_PROVIDER == "none":
    EMBED_MODEL = ""
    EMBED_DIM   = 0
    EMBED_BATCH = 0
elif EMBED_PROVIDER == "openai":
    EMBED_MODEL = "text-embedding-3-small"
    EMBED_DIM   = 1536
    EMBED_BATCH = 64
elif EMBED_PROVIDER == "gemini":
    # gemini-embedding-001 supports outputDimensionality truncation; we pin
    # 768 to keep storage compact. (text-embedding-004 was deprecated.)
    EMBED_MODEL = "gemini-embedding-001"
    EMBED_DIM   = 768
    EMBED_BATCH = 1
else:
    # Self-hosted ONNX via fastembed — zero per-call cost, no API key.
    # paraphrase-multilingual-MiniLM-L12-v2: 384-dim, ~120MB ONNX,
    # supports 50+ languages including strong Korean + English.
    EMBED_PROVIDER = "local"
    EMBED_MODEL = "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
    EMBED_DIM   = 384
    EMBED_BATCH = 32
# Rough char→token ratio for English ≈ 4; Korean ≈ 2-3. Use 800 chars ≈ 200-300 tokens
# per chunk so we stay well below the 8192-token embedding limit and keep semantic
# units coherent.
MAX_CHARS_PER_CHUNK = 1200
ROWS_PER_BLOCK      = 8     # xlsx/csv: group this many data rows together


# ---------------------------------------------------------------------------
# Parsers — each returns list[dict(location, title, content)]
# ---------------------------------------------------------------------------

def _parse_xlsx(filename: str, blob: bytes) -> list[dict]:
    """Split each sheet into N-row blocks; prepend column headers so each
    chunk is self-describing ('회사=트리플에이치 | 구분=상가 | ...').
    Header detection: the first non-empty row whose ≥60% of cells are strings
    is treated as the header row; rows above it (often dates / titles) are
    captured as a 'preface' chunk per sheet."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(blob), read_only=True, data_only=True)
    chunks: list[dict] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        # Trim trailing empty rows
        while rows and all(c is None or (isinstance(c, str) and not c.strip()) for c in rows[-1]):
            rows.pop()
        if not rows:
            continue

        # Find header row: first row with ≥2 cells AND ≥60% string-type cells
        header_idx = -1
        for i, r in enumerate(rows[:10]):
            non_empty = [c for c in r if c is not None and str(c).strip()]
            if len(non_empty) < 2:
                continue
            str_cells = sum(1 for c in non_empty if isinstance(c, str))
            if str_cells / len(non_empty) >= 0.6:
                header_idx = i
                break

        if header_idx >= 0:
            headers = [str(c).strip() if c is not None else f"col{j+1}" for j, c in enumerate(rows[header_idx])]
            preface_rows = rows[:header_idx]
            data_rows    = rows[header_idx + 1:]
        else:
            headers = [f"col{j+1}" for j in range(max((len(r) for r in rows), default=0))]
            preface_rows = []
            data_rows    = rows

        # Preface (titles / dates above the header) — one chunk
        if preface_rows:
            text = "\n".join(
                " | ".join(str(c) for c in r if c is not None and str(c).strip())
                for r in preface_rows
            )
            if text.strip():
                chunks.append({
                    "location": f"Sheet: {sheet_name} / preface",
                    "title": sheet_name,
                    "content": f"[Sheet: {sheet_name}]\n{text}",
                })

        # Data rows in blocks of ROWS_PER_BLOCK
        block: list[str] = []
        block_start = header_idx + 2 if header_idx >= 0 else 1  # 1-based
        for i, r in enumerate(data_rows):
            non_empty = [c for c in r if c is not None and str(c).strip()]
            if not non_empty:
                continue
            pairs = []
            for j, cell in enumerate(r):
                if cell is None:
                    continue
                v = str(cell).strip()
                if not v:
                    continue
                h = headers[j] if j < len(headers) else f"col{j+1}"
                pairs.append(f"{h}={v}")
            block.append(" | ".join(pairs))
            if len(block) >= ROWS_PER_BLOCK:
                block_end = block_start + len(block) - 1
                content = "\n".join(block)
                chunks.append({
                    "location": f"Sheet: {sheet_name} / rows {block_start}-{block_end}",
                    "title": sheet_name,
                    "content": f"[Sheet: {sheet_name}]\nHeaders: {' | '.join(headers)}\n{content}",
                })
                block_start = block_end + 1
                block = []
        if block:
            block_end = block_start + len(block) - 1
            content = "\n".join(block)
            chunks.append({
                "location": f"Sheet: {sheet_name} / rows {block_start}-{block_end}",
                "title": sheet_name,
                "content": f"[Sheet: {sheet_name}]\nHeaders: {' | '.join(headers)}\n{content}",
            })

    wb.close()
    return chunks


def _parse_pdf(filename: str, blob: bytes) -> list[dict]:
    """One chunk per page. If a page is huge (>MAX_CHARS_PER_CHUNK * 3) split it."""
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader  # type: ignore
    reader = PdfReader(io.BytesIO(blob))
    out: list[dict] = []
    for i, page in enumerate(reader.pages):
        try:
            txt = page.extract_text() or ""
        except Exception as e:
            log.warning("pdf page %d extract failed: %s", i, e)
            txt = ""
        txt = txt.strip()
        if not txt:
            continue
        # Big pages → sub-split
        pieces = _split_text(txt, MAX_CHARS_PER_CHUNK)
        for k, piece in enumerate(pieces):
            loc = f"page {i+1}" + (f" / part {k+1}" if len(pieces) > 1 else "")
            out.append({"location": loc, "title": f"{filename} — p.{i+1}", "content": piece})
    return out


def _parse_docx(filename: str, blob: bytes) -> list[dict]:
    """Paragraph-level chunks, accumulated up to MAX_CHARS_PER_CHUNK."""
    from docx import Document
    doc = Document(io.BytesIO(blob))
    out: list[dict] = []
    buf: list[str] = []
    para_start = 1
    cur = 0
    last_heading: Optional[str] = None
    for idx, p in enumerate(doc.paragraphs, start=1):
        txt = (p.text or "").strip()
        if not txt:
            continue
        if (p.style and p.style.name and p.style.name.lower().startswith("heading")):
            last_heading = txt
        if cur + len(txt) > MAX_CHARS_PER_CHUNK and buf:
            out.append({
                "location": f"paragraphs {para_start}-{idx-1}",
                "title": last_heading or filename,
                "content": "\n".join(buf),
            })
            buf = []
            cur = 0
            para_start = idx
        buf.append(txt)
        cur += len(txt) + 1
    if buf:
        out.append({
            "location": f"paragraphs {para_start}-{len(doc.paragraphs)}",
            "title": last_heading or filename,
            "content": "\n".join(buf),
        })
    return out


def _parse_pptx(filename: str, blob: bytes) -> list[dict]:
    """One chunk per slide, including title + body text + notes."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    out: list[dict] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                if not title and shape == slide.shapes.title if slide.shapes.title else False:
                    title = line
                else:
                    body_parts.append(line)
        # Slide notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        text = "\n".join(body_parts)
        if notes:
            text = f"{text}\n\n[Notes]\n{notes}"
        text = text.strip()
        if not text and not title:
            continue
        out.append({
            "location": f"slide {i}",
            "title": title or f"Slide {i}",
            "content": f"{title}\n{text}" if title else text,
        })
    return out


def _parse_csv(filename: str, blob: bytes) -> list[dict]:
    import csv as _csv
    text_in = blob.decode("utf-8-sig", errors="replace")
    rdr = list(_csv.reader(io.StringIO(text_in)))
    if not rdr:
        return []
    headers = [str(c).strip() or f"col{i+1}" for i, c in enumerate(rdr[0])]
    out: list[dict] = []
    block: list[str] = []
    start = 2
    for i, r in enumerate(rdr[1:], start=2):
        pairs = []
        for j, cell in enumerate(r):
            v = (cell or "").strip()
            if not v:
                continue
            h = headers[j] if j < len(headers) else f"col{j+1}"
            pairs.append(f"{h}={v}")
        if pairs:
            block.append(" | ".join(pairs))
        if len(block) >= ROWS_PER_BLOCK:
            end = start + len(block) - 1
            out.append({
                "location": f"rows {start}-{end}",
                "title": filename,
                "content": f"Headers: {' | '.join(headers)}\n" + "\n".join(block),
            })
            start = end + 1
            block = []
    if block:
        end = start + len(block) - 1
        out.append({
            "location": f"rows {start}-{end}",
            "title": filename,
            "content": f"Headers: {' | '.join(headers)}\n" + "\n".join(block),
        })
    return out


def _parse_text(filename: str, blob: bytes) -> list[dict]:
    text = blob.decode("utf-8", errors="replace")
    pieces = _split_text(text, MAX_CHARS_PER_CHUNK)
    out: list[dict] = []
    for i, piece in enumerate(pieces):
        out.append({
            "location": f"part {i+1}/{len(pieces)}",
            "title": filename,
            "content": piece,
        })
    return out


def _split_text(text: str, max_chars: int) -> list[str]:
    """Split text into chunks at paragraph boundaries when possible."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]
    paras = re.split(r"\n\s*\n", text)
    out: list[str] = []
    cur = ""
    for p in paras:
        p = p.strip()
        if not p:
            continue
        if len(p) > max_chars:
            # hard-split very long paragraph
            for i in range(0, len(p), max_chars):
                out.append(p[i:i + max_chars])
            continue
        if len(cur) + len(p) + 2 > max_chars:
            if cur:
                out.append(cur)
            cur = p
        else:
            cur = (cur + "\n\n" + p) if cur else p
    if cur:
        out.append(cur)
    return out


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_EXT_TO_PARSER = {
    ".xlsx": _parse_xlsx,
    ".xlsm": _parse_xlsx,
    ".xls":  _parse_xlsx,    # openpyxl reads xlsx; old xls may fail — caller catches
    ".csv":  _parse_csv,
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt":  _parse_text,
    ".md":   _parse_text,
    ".json": _parse_text,
}

SUPPORTED_EXTS = sorted(_EXT_TO_PARSER.keys())


def parse_file(filename: str, blob: bytes) -> list[dict]:
    """Parse a file into chunks. Raises ValueError on unsupported extension."""
    name = filename.lower()
    for ext, parser in _EXT_TO_PARSER.items():
        if name.endswith(ext):
            return parser(filename, blob)
    raise ValueError(
        f"Unsupported file type for '{filename}'. Supported: {', '.join(SUPPORTED_EXTS)}"
    )


# ---------------------------------------------------------------------------
# Embedding
# ---------------------------------------------------------------------------

def _openai_key() -> Optional[str]:
    return os.getenv("OPENAI_API_KEY") or os.getenv("LLM_API_KEY")


def _gemini_key() -> Optional[str]:
    return (
        os.getenv("GOOGLE_GENERATIVE_AI_API_KEY")
        or os.getenv("GEMINI_API_KEY")
        or os.getenv("GOOGLE_API_KEY")
    )


def _embed_openai(safe: list[str]) -> list[list[float]]:
    key = _openai_key()
    if not key:
        raise RuntimeError("OPENAI_API_KEY is not set; cannot generate embeddings")
    resp = httpx.post(
        "https://api.openai.com/v1/embeddings",
        headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
        json={"model": EMBED_MODEL, "input": safe},
        timeout=60.0,
    )
    if resp.status_code != 200:
        raise RuntimeError(f"OpenAI embed failed {resp.status_code}: {resp.text[:300]}")
    return [row["embedding"] for row in resp.json()["data"]]


def _embed_gemini(safe: list[str]) -> list[list[float]]:
    """Embed via Gemini's gemini-embedding-001 (3072-dim native, truncated to
    EMBED_DIM via outputDimensionality). The synchronous v1beta surface only
    exposes per-text embedContent; we loop here. For our ~700-chunk workbook
    that's ~30s — acceptable for a one-time ingest."""
    key = _gemini_key()
    if not key:
        raise RuntimeError(
            "GEMINI_API_KEY / GOOGLE_GENERATIVE_AI_API_KEY is not set; "
            "cannot generate embeddings"
        )
    out: list[list[float]] = []
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{EMBED_MODEL}:embedContent"
    with httpx.Client(timeout=60.0) as client:
        for t in safe:
            body = {
                "model": f"models/{EMBED_MODEL}",
                "content": {"parts": [{"text": t}]},
                "outputDimensionality": EMBED_DIM,
            }
            resp = client.post(url, params={"key": key}, json=body)
            if resp.status_code != 200:
                raise RuntimeError(
                    f"Gemini embed failed {resp.status_code}: {resp.text[:300]}"
                )
            data = resp.json()
            vec = (data.get("embedding") or {}).get("values") or []
            # gemini-embedding-001 returns unit-normalised vectors only when
            # outputDimensionality == native (3072). For truncated dims we
            # re-normalise so cosine similarity stays meaningful.
            if EMBED_DIM != 3072 and vec:
                import math
                norm = math.sqrt(sum(v * v for v in vec)) or 1.0
                vec = [v / norm for v in vec]
            out.append(vec)
    return out


# Lazy-loaded fastembed model singleton — first call downloads ~120MB ONNX
# into the user's HF cache; subsequent calls reuse it instantly.
_local_model = None

def _embed_local(safe: list[str]) -> list[list[float]]:
    """Embed via on-disk ONNX model (intfloat/multilingual-e5-small).
    The e5 family requires query/passage prefixes for best quality. For our
    'asymmetric' use case (ingest = passage, query = query) we prefix
    accordingly inside ingest_file and rag_retrieve."""
    global _local_model
    if _local_model is None:
        try:
            from fastembed import TextEmbedding
        except ImportError as e:
            raise RuntimeError(
                "fastembed not installed. Run: pip install fastembed"
            ) from e
        _local_model = TextEmbedding(model_name=EMBED_MODEL)
    # fastembed yields numpy arrays; cast to Python lists
    return [list(map(float, v)) for v in _local_model.embed(safe)]


def embed_batch(texts: list[str]) -> list[list[float]]:
    """Synchronously embed up to EMBED_BATCH texts via the configured provider.
    Returns one vector per input. Raises RuntimeError on auth / network failures."""
    if not texts:
        return []
    safe = [t if (t and t.strip()) else " " for t in texts]
    if EMBED_PROVIDER == "openai":
        return _embed_openai(safe)
    if EMBED_PROVIDER == "gemini":
        return _embed_gemini(safe)
    return _embed_local(safe)


def _vec_to_pg(v: list[float]) -> str:
    """Render a vector list as pgvector's text format: '[0.1,0.2,...]'."""
    return "[" + ",".join(f"{x:.6f}" for x in v) + "]"


# ---------------------------------------------------------------------------
# Ingestion
# ---------------------------------------------------------------------------

def ingest_file(
    db: Session,
    *,
    agent_id: str,
    filename: str,
    mime_type: Optional[str],
    blob: bytes,
    uploaded_by: Optional[str] = None,
) -> dict:
    """Parse + embed + insert. Returns {file_id, chunk_count, status}.
    Caller is responsible for db.commit()."""
    # 1. Insert file row (status=pending)
    file_id = str(uuid.uuid4())
    db.execute(sa_text("""
        INSERT INTO assistant_knowledge_files
            (id, agent_id, filename, mime_type, size_bytes, uploaded_by, status, chunk_count)
        VALUES (:id, :agent_id, :filename, :mime_type, :size_bytes, :uploaded_by, 'pending', 0)
    """), {
        "id": file_id,
        "agent_id": agent_id,
        "filename": filename,
        "mime_type": mime_type,
        "size_bytes": len(blob),
        "uploaded_by": uploaded_by,
    })
    db.commit()

    try:
        chunks = parse_file(filename, blob)
    except Exception as e:
        db.execute(sa_text("""
            UPDATE assistant_knowledge_files
               SET status='error', error_msg=:msg
             WHERE id=:id
        """), {"id": file_id, "msg": f"parse failed: {e}"[:1000]})
        db.commit()
        raise

    if not chunks:
        db.execute(sa_text("""
            UPDATE assistant_knowledge_files
               SET status='indexed', chunk_count=0, error_msg='no parseable content'
             WHERE id=:id
        """), {"id": file_id})
        db.commit()
        return {"file_id": file_id, "chunk_count": 0, "status": "empty"}

    # 2. Insert chunks. When EMBED_PROVIDER=none we skip the embedding column;
    # retrieval still works via pgroonga full-text-search on `content`.
    inserted = 0
    use_embeddings = EMBED_PROVIDER != "none"
    passage_prefix = "passage: " if (EMBED_PROVIDER == "local" and "e5" in EMBED_MODEL.lower()) else ""
    try:
        batch_size = EMBED_BATCH if use_embeddings else 64
        for start in range(0, len(chunks), batch_size):
            batch = chunks[start:start + batch_size]
            if use_embeddings:
                texts = [(passage_prefix + c["content"]) for c in batch]
                vecs  = embed_batch(texts)
            else:
                vecs = [None] * len(batch)
            for c, v in zip(batch, vecs):
                if use_embeddings and v is not None:
                    db.execute(sa_text("""
                        INSERT INTO assistant_knowledge_chunks
                            (file_id, agent_id, location, title, content, token_count,
                             embedding, metadata)
                        VALUES
                            (:file_id, :agent_id, :location, :title, :content, :tokens,
                             CAST(:embedding AS vector), CAST(:metadata AS jsonb))
                    """), {
                        "file_id":   file_id,
                        "agent_id":  agent_id,
                        "location":  c.get("location"),
                        "title":     c.get("title"),
                        "content":   c["content"],
                        "tokens":    len(c["content"]) // 4,
                        "embedding": _vec_to_pg(v),
                        "metadata":  json.dumps(c.get("metadata") or {}),
                    })
                else:
                    db.execute(sa_text("""
                        INSERT INTO assistant_knowledge_chunks
                            (file_id, agent_id, location, title, content, token_count, metadata)
                        VALUES
                            (:file_id, :agent_id, :location, :title, :content, :tokens,
                             CAST(:metadata AS jsonb))
                    """), {
                        "file_id":   file_id,
                        "agent_id":  agent_id,
                        "location":  c.get("location"),
                        "title":     c.get("title"),
                        "content":   c["content"],
                        "tokens":    len(c["content"]) // 4,
                        "metadata":  json.dumps(c.get("metadata") or {}),
                    })
                inserted += 1
            db.commit()
    except Exception as e:
        db.execute(sa_text("""
            UPDATE assistant_knowledge_files
               SET status='error', chunk_count=:n, error_msg=:msg
             WHERE id=:id
        """), {"id": file_id, "n": inserted, "msg": f"ingest failed: {e}"[:1000]})
        db.commit()
        raise

    db.execute(sa_text("""
        UPDATE assistant_knowledge_files
           SET status='indexed', chunk_count=:n
         WHERE id=:id
    """), {"id": file_id, "n": inserted})
    db.commit()
    return {"file_id": file_id, "chunk_count": inserted, "status": "indexed"}


# ---------------------------------------------------------------------------
# Retrieval
# ---------------------------------------------------------------------------

def rag_retrieve(
    db: Session,
    *,
    agent_id: str,
    query: str,
    top_k: int = 8,
    min_sim: float = 0.35,
) -> list[dict]:
    """Vector-search the agent's KB for `query`. Returns [] when nothing scores
    above `min_sim` — caller then decides whether to fall back to the LLM's
    own knowledge. Each dict has: content, location, title, similarity, file_id, filename."""
    q = (query or "").strip()
    if not q:
        return []

    # --- Text-only path (pgroonga, no embedding API needed) ---
    if EMBED_PROVIDER == "none":
        try:
            rows = db.execute(sa_text("""
                SELECT * FROM search_assistant_knowledge_text(
                    :agent_id, :query, :top_k
                )
            """), {
                "agent_id": agent_id,
                "query":    q,
                "top_k":    top_k,
            }).fetchall()
        except Exception as e:
            log.warning("pgroonga retrieve failed (%s) — returning empty", e)
            return []
        return [
            {
                "chunk_id":   str(r.chunk_id),
                "file_id":    str(r.file_id),
                "filename":   r.filename,
                "location":   r.location,
                "title":      r.title,
                "content":    r.content,
                "similarity": float(r.similarity),
            }
            for r in rows
        ]

    # --- Vector path (when embeddings are configured) ---
    q_input = ("query: " + q) if (EMBED_PROVIDER == "local" and "e5" in EMBED_MODEL.lower()) else q
    try:
        vec = embed_batch([q_input])[0]
    except Exception as e:
        log.warning("rag_retrieve: embedding failed (%s) — returning empty", e)
        return []

    rows = db.execute(sa_text("""
        SELECT * FROM search_assistant_knowledge(
            :agent_id, CAST(:embedding AS vector), :top_k, :min_sim
        )
    """), {
        "agent_id":  agent_id,
        "embedding": _vec_to_pg(vec),
        "top_k":     top_k,
        "min_sim":   min_sim,
    }).fetchall()

    return [
        {
            "chunk_id":   str(r.chunk_id),
            "file_id":    str(r.file_id),
            "filename":   r.filename,
            "location":   r.location,
            "title":      r.title,
            "content":    r.content,
            "similarity": float(r.similarity),
        }
        for r in rows
    ]


def list_files(db: Session, *, agent_id: str) -> list[dict]:
    rows = db.execute(sa_text("""
        SELECT id, filename, mime_type, size_bytes, uploaded_by, uploaded_at,
               status, chunk_count, error_msg
          FROM assistant_knowledge_files
         WHERE agent_id = :agent_id
         ORDER BY uploaded_at DESC
    """), {"agent_id": agent_id}).fetchall()
    return [
        {
            "id":          str(r.id),
            "filename":    r.filename,
            "mime_type":   r.mime_type,
            "size_bytes":  r.size_bytes,
            "uploaded_by": r.uploaded_by,
            "uploaded_at": r.uploaded_at.isoformat() if r.uploaded_at else None,
            "status":      r.status,
            "chunk_count": r.chunk_count,
            "error_msg":   r.error_msg,
        } for r in rows
    ]


def delete_file(db: Session, *, agent_id: str, file_id: str) -> int:
    """Returns number of file rows deleted (0 or 1). Chunks cascade."""
    res = db.execute(sa_text("""
        DELETE FROM assistant_knowledge_files
         WHERE id=:id AND agent_id=:agent_id
    """), {"id": file_id, "agent_id": agent_id})
    db.commit()
    return res.rowcount or 0
